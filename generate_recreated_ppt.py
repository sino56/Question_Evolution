from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "Question_Evolution_流程治理与评分闭环.pptx"

SLIDE_W = 10.0
SLIDE_H = 5.625
PX_W = 592
PX_H = 333


def X(px: float) -> float:
    return px * SLIDE_W / PX_W


def Y(px: float) -> float:
    return px * SLIDE_H / PX_H


FONT = "Microsoft YaHei"
FONT_SCALE = 1.22
BLACK = "222222"
MUTED = "666C70"
LINE = "9AA6A7"
PALE_LINE = "BEC4C4"
PINK = "FEDCDD"
BLUE = "CFE6FF"
MINT = "CCF5DC"
YELLOW = "FFF0B2"
GRAY = "F0F5F8"
RED = "A92929"
BLUE_TEXT = "245EAB"
GREEN_TEXT = "198C6C"
AMBER = "9B5A13"


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.strip("#")
    return RGBColor.from_string(value)


def set_shape_transparency(shape, alpha: int = 0) -> None:
    solid = shape.fill._xPr.solidFill
    if solid is None:
        return
    srgb = solid.find("{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
    if srgb is None:
        return
    alpha_node = OxmlElement("a:alpha")
    alpha_node.set("val", str(max(0, min(100000, (100 - alpha) * 1000))))
    srgb.append(alpha_node)


def set_cell_margins(tf, left=4, right=4, top=2, bottom=2):
    tf.margin_left = Pt(left)
    tf.margin_right = Pt(right)
    tf.margin_top = Pt(top)
    tf.margin_bottom = Pt(bottom)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 10,
    color: str = BLACK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    font: str = FONT,
    margin: float = 0,
    rotation: float = 0,
    line_spacing: float | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.rotation = rotation
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = valign
    set_cell_margins(tf, margin, margin, margin, margin)
    parts = text.split("\n")
    for i, part in enumerate(parts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = part
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        if line_spacing is not None:
            p.line_spacing = line_spacing
        p.font.name = font
        p.font.size = Pt(size * FONT_SCALE)
        p.font.bold = bold
        p.font.color.rgb = rgb(color)
    return box


def add_round_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str,
    line: str = PALE_LINE,
    line_width: float = 0.75,
    radius: bool = True,
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(line_width)
    return shape


def add_box_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = GRAY,
    line: str = PALE_LINE,
    size: float = 9.5,
    color: str = BLACK,
    bold: bool = True,
    align: PP_ALIGN = PP_ALIGN.CENTER,
    valign: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
    margin: float = 3,
    radius: bool = True,
):
    shape = add_round_rect(slide, x, y, w, h, fill=fill, line=line, radius=radius)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = valign
    set_cell_margins(tf, margin, margin, margin, margin)
    for i, part in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = part
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.font.name = FONT
        p.font.size = Pt(size * FONT_SCALE)
        p.font.bold = bold
        p.font.color.rgb = rgb(color)
    return shape


def add_card(
    slide,
    title: str,
    lines: list[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str,
    title_color: str,
    body_size: float = 9.3,
    title_size: float = 10.5,
    bullets: bool = True,
):
    shape = add_round_rect(slide, x, y, w, h, fill=fill, line=PALE_LINE, line_width=0.8)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = MSO_ANCHOR.TOP
    set_cell_margins(tf, 8, 7, 6, 4)
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.size = Pt(title_size * FONT_SCALE)
    p.font.bold = True
    p.font.color.rgb = rgb(title_color)
    p.space_after = Pt(2)
    for value in lines:
        p = tf.add_paragraph()
        p.text = ("• " if bullets else "") + value
        p.font.name = FONT
        p.font.size = Pt(body_size * FONT_SCALE)
        p.font.bold = True
        p.font.color.rgb = rgb(BLACK)
        p.space_before = Pt(0)
        p.space_after = Pt(1)
        p.line_spacing = 1.0
    return shape


def add_line(slide, x1, y1, x2, y2, *, color=LINE, width=1.0, arrow=False):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if arrow:
        ln = line.line._get_or_add_ln()
        tail = OxmlElement("a:tailEnd")
        tail.set("type", "triangle")
        tail.set("w", "sm")
        tail.set("len", "sm")
        ln.append(tail)
    return line


def add_header(slide, title: str, subtitle: str, *, title_y=5, sub_y=25, line_y=38):
    add_text(slide, title, X(6), Y(title_y), X(550), Y(20), size=14.2, bold=True)
    add_text(slide, subtitle, X(7), Y(sub_y), X(548), Y(12), size=6.7, color=MUTED, bold=True)
    add_line(slide, X(7), Y(line_y), X(549), Y(line_y), color=LINE, width=1.3)


def add_section(slide, title: str, subtitle: str, *, title_y: int, sub_y: int, line_y: int):
    add_text(slide, title, X(6), Y(title_y), X(555), Y(22), size=13.1, bold=True)
    add_text(slide, subtitle, X(7), Y(sub_y), X(555), Y(12), size=6.5, color=MUTED, bold=True)
    add_line(slide, X(7), Y(line_y), X(549), Y(line_y), color=LINE, width=1.2)


def add_page_number(slide, number: int):
    add_box_text(
        slide,
        str(number),
        X(8),
        Y(303),
        X(35),
        Y(23),
        fill="333333",
        line="222222",
        color="FFFFFF",
        size=14,
        bold=False,
        margin=0,
        radius=True,
    )


def build_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "进化过程拆解分析", "同一个进化决策可以垂直重置题面内容，也可按流程自然排除源，必须分层定位。")
    add_card(slide, "未拆分时的风险", ["题面暗示答案方向", "旧参考答案污染新题", "路由 fallback 掩盖空选", "校验误挡在线程阻断"], X(13), Y(46), X(166), Y(96), fill=PINK, title_color=RED)
    add_card(slide, "拆分后的职责", ["内容侧：定义算子题面边界", "流程侧：执行授权和路由门禁", "评分侧：重建答案与 Rubric", "经验侧：只写可验证结论"], X(195), Y(46), X(166), Y(96), fill=BLUE, title_color=BLUE_TEXT)
    add_card(slide, "判断标准", ["能回答", "不泄漏", "可评分", "能归因", "可复盘"], X(377), Y(46), X(145), Y(96), fill=MINT, title_color=GREEN_TEXT, body_size=9.0)

    add_section(slide, "总设计：两条控制链", "内容侧管“公开题面”，流程机制管“运行闭环”，二者通过授权和诊断结果连接。", title_y=157, sub_y=176, line_y=187)
    add_box_text(slide, "算子内容治理\n规定题面内容 + 执行契约 + 内容测试", X(11), Y(196), X(244), Y(47), fill=BLUE, line="C0CAD0", size=9.4, color=BLUE_TEXT)
    add_box_text(slide, "流程机制治理\n授权 + 模式 + 路由 + 校验 + 评分闭环", X(291), Y(196), X(242), Y(47), fill=MINT, line="B4CBBE", size=9.4, color=GREEN_TEXT)
    add_box_text(slide, "中性题面\n禁止答案提示", X(24), Y(251), X(71), Y(45), fill="FFFFFF", size=7.6, margin=2)
    add_box_text(slide, "槽位契约\n声明完成题目生成\n或需要材料", X(166), Y(251), X(78), Y(45), fill="FFFFFF", size=7.0, margin=1)
    add_box_text(slide, "弱模型\n只看最终题面", X(304), Y(251), X(73), Y(45), fill="FFFFFF", size=7.6, margin=1)
    add_box_text(slide, "新题\n重建答案与评分", X(389), Y(251), X(73), Y(45), fill="FFFFFF", size=7.6, margin=1)
    add_box_text(slide, "效果\n真实得分才写\n经验", X(475), Y(251), X(69), Y(45), fill="FFFFFF", size=7.1, margin=1)
    add_box_text(slide, "共同目标：有效进化可验证、可追溯、可迭代", X(170), Y(305), X(229), Y(22), fill=YELLOW, line="E3CF8A", size=8.6, color=AMBER, margin=1)
    add_page_number(slide, 1)


def build_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "内容治理原理：题面只给可观察材料", "因为题目能力目标应落在一个中性业务判断任务，内题界面提供材料而不是提前暗示。")
    add_card(slide, "题面可以写", ["可观察事实：时间、地点、对象、数值", "题内明示规则、公式、单位、局部参数", "竞争解释所需的对称事实", "一个用户可见的业务判断任务"], X(6), Y(43), X(247), Y(101), fill=MINT, title_color=GREEN_TEXT, body_size=8.8)
    add_card(slide, "题面不能写", ["答案方向：不能认定、最高只能支持", "事实角色：关键证据、干扰事实", "推理链总结：哪一跳成立或缺失", "规则应用结果和竞争解释淘汰"], X(303), Y(43), X(247), Y(101), fill=PINK, title_color=RED, body_size=8.8)
    add_text(slide, "核心原则：公开事实足够答题，但不能提前完成答案归纳。", X(166), Y(151), X(265), Y(16), size=8.3, color=BLUE_TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_section(slide, "执行契约：把算子能力转成可检查槽位", "每个生成算子都用契约说明适用模式、材料槽位、可合成边界和必然控制。", title_y=174, sub_y=193, line_y=204)
    flow_x = [31, 142, 253, 372]
    flow_w = [89, 90, 93, 116]
    flow_t = ["算子能力轴", "算子必须具备的材\n料信息", "合理的补充材料信\n息", "不能输送的材料信息"]
    for i, (x, w, t) in enumerate(zip(flow_x, flow_w, flow_t)):
        add_box_text(slide, t, X(x), Y(217), X(w), Y(35), fill=BLUE, line="C2CDD2", size=7.5, margin=1)
        if i < 3:
            add_line(slide, X(x + w + 4), Y(234), X(flow_x[i + 1] - 4), Y(234), color=BLUE_TEXT, width=1.0, arrow=True)
    add_text(slide, "必须来自来源、规则或真实依据", X(359), Y(254), X(178), Y(12), size=6.5, bold=True, align=PP_ALIGN.CENTER)
    add_box_text(slide, "必需材料\n路径、时间窗、锚点观测、规则\n字段等。\n缺失时记录槽位不足。", X(12), Y(270), X(151), Y(61), fill="FFFFFF", size=7.4, margin=1)
    add_box_text(slide, "可题内合成\n人物、实体、动作、路径、局部\n参数、竞争解释。\n必须登记同世界事实。", X(178), Y(270), X(158), Y(61), fill="FFFFFF", size=7.4, margin=1)
    add_box_text(slide, "不可合成\n真实法规、真实专业阈值、真实\n外部记录、真实案件结论。\n只能来自可审计来源。", X(349), Y(270), X(198), Y(61), fill="FFFFFF", size=7.4, margin=1)
    add_page_number(slide, 2)


def build_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "流程机制原理：正式实验必须跑完整闭环", "信任题面和流程中的“任务力”，真正按新题重建答案、重新评分、每轮回到效果归因。")
    labels = ["重建参考答案", "重建 Rubric", "弱模型回答", "正式评分", "效果归因", "下一轮"]
    xs = [24, 116, 211, 319, 408, 501]
    ws = [79, 86, 88, 72, 75, 49]
    for i, (label, x, w) in enumerate(zip(labels, xs, ws)):
        add_box_text(slide, label, X(x), Y(48), X(w), Y(30), fill=MINT, line="B9D5C7", size=7.8, margin=1)
        if i < len(labels) - 1:
            add_line(slide, X(x + w + 4), Y(63), X(xs[i + 1] - 4), Y(63), color=GREEN_TEXT, width=1.0, arrow=True)
    add_card(slide, "正式入口红线", ["缺少已配不能跳过阶段 → 直接报错", "旧材料不能硬塞进活动评分 → 回到方向审核材料"], X(11), Y(88), X(239), Y(54), fill=YELLOW, title_color=AMBER, body_size=7.4, title_size=9.5)
    add_card(slide, "校验职责边界", ["默认记录诊断和建议动作 → 记录风险以及是否人工审核", "只有技术不可执行才停止链路 → 格式无法解析"], X(299), Y(88), X(254), Y(54), fill=GRAY, title_color=BLUE_TEXT, body_size=7.2, title_size=9.2)
    add_section(slide, "材料隔离：弱模型只看最终题面", "题面变化后，旧参考答案、旧评分标准、预测答案和隐含链路全部成为隔离材料。", title_y=151, sub_y=170, line_y=182)
    add_box_text(slide, "旧材料\n参考答案 / Rubric / 分数 / 旧题\n失败", X(26), Y(191), X(145), Y(52), fill=PINK, color=RED, size=8.0, margin=1)
    add_line(slide, X(176), Y(217), X(195), Y(217), color=GREEN_TEXT, arrow=True)
    add_box_text(slide, "隔离层\n标记过期，不进入弱模型上下\n文", X(206), Y(191), X(145), Y(52), fill=YELLOW, color=AMBER, size=8.0, margin=1)
    add_line(slide, X(357), Y(217), X(377), Y(217), color=GREEN_TEXT, arrow=True)
    add_box_text(slide, "弱模型\n只接收 evolved_prompt", X(385), Y(191), X(146), Y(52), fill=MINT, color=GREEN_TEXT, size=8.0, margin=1)
    add_card(slide, "来源拆分", ["观察事实、题目本身任务、题中隐含规则", "答案方向和推导总结不得进入公开题面", "事实版本记录：事实的编号 / 属于哪个场景 / 信息来源"], X(31), Y(255), X(241), Y(62), fill="FFFFFF", title_color=BLUE_TEXT, body_size=7.2, title_size=9.0)
    add_card(slide, "评分状态", ["新题活动评分状态标记为“还没正式评分”", "真实评分前只能说结构有效，不能说是否有效进化", "有效性必须由新评分闭环证明"], X(294), Y(255), X(240), Y(62), fill="FFFFFF", title_color=GREEN_TEXT, body_size=7.2, title_size=9.0)
    add_page_number(slide, 3)


def build_slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "模式决策：先分析，再判断如何进化", "模式决策器输入来源分析、样本画像、任务授权和允许模式，输出可审计的 authorization_id（这道题允许这样改，是根据哪些配置或任务授权来的）。")
    rows = [
        ("只忠实使用原题来源", "新题只能基于原题已有事实和规则重组、改问、强调某个推理点，不能凭空补新事实"),
        ("构造一个明确的题内假设案例", "如果原题太抽象或事实太少，可以新建一个“本题设定”的小案例，用来考目标能力"),
        ("在原题基础上受控补材料", "保留原题核心事实，同时为了形成目标推理结构，补一些题内设定、人物、路径、局部参数、竞争解释"),
        ("从原题改编出假设版本", "以原题为基准，改造成一个假设场景。它保留原题的结构或核心关系，但允许替换实体、数值、路径、场景，用来测试同类能力。"),
        ("不安全进化，直接透传原题", "当前样本没有安全可用的进化方式，或者生成的新题技术上不可执行，\n就不强行改题，直接保留原题继续流程或记录原因。"),
    ]
    ys = [54, 94, 134, 174, 214]
    for i, ((label, text), y) in enumerate(zip(rows, ys)):
        fill = YELLOW if i == 4 else GRAY
        color = AMBER if i == 4 else BLACK
        add_box_text(slide, label, X(8), Y(y), X(136), Y(27), fill=fill, line="D5D9D9", size=8.0, color=color, margin=1)
        add_text(slide, text, X(155), Y(y + 2), X(400), Y(35 if i >= 2 else 26), size=8.2, bold=True, valign=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    add_box_text(slide, "原则：授权限制真实外部事实与真实规则伪造，不禁止明确标注的题内假设构造。", X(37), Y(263), X(484), Y(25), fill=BLUE, line="AFC6D6", size=8.1, color=BLUE_TEXT, margin=1)
    add_page_number(slide, 4)


def build_slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "模式感知路由：关键词只召回候选，不携带风格", "路由进入命题链路后只做排序，按授权范围和经验决定当前候选。")
    top = [
        ("任务授权信息", "任务允许什么"),
        ("原化模式摘要", "适用模式要素\n支持哪些算子"),
        ("重点槽位检查", "需要的材料是否满足"),
        ("推理连环筛选", "选择最小开放"),
        ("经验偏好序", "参考历史经验"),
    ]
    xs = [69, 157, 245, 331, 414]
    ws = [73, 77, 75, 70, 57]
    for (label, desc), x, w in zip(top, xs, ws):
        add_box_text(slide, label, X(x), Y(48), X(w), Y(30), fill=BLUE, line="C3CED4", size=7.3, margin=1)
        add_text(slide, desc, X(x - 7), Y(82), X(w + 14), Y(25), size=6.8, bold=True, align=PP_ALIGN.CENTER)
    add_section(slide, "隐藏规划与盲写题面：上下文隔离", "同一次模型调用不应既暴露策略规划又要求盲写公开题面，必须分出全量规划与公开事实投影。", title_y=115, sub_y=134, line_y=145)
    add_box_text(slide, "问题规划器\n\n事实内容\n命题意图\n证据覆盖顺序\n推理结构\n控制计划\n预测弱模型错误", X(7), Y(153), X(228), Y(126), fill=BLUE, line="B6C8D6", size=8.6, color=BLUE_TEXT, margin=2)
    add_line(slide, X(238), Y(217), X(252), Y(217), color=BLUE_TEXT, arrow=True)
    add_box_text(slide, "公开事实投影", X(254), Y(198), X(80), Y(38), fill="FFFFFF", line="C6CBCB", size=8.0, margin=1)
    add_text(slide, "把隐含规划器里的材料和过滤\n一遍，只留下会题可见信息", X(245), Y(241), X(102), Y(30), size=6.5, bold=True, align=PP_ALIGN.CENTER)
    add_line(slide, X(339), Y(217), X(354), Y(217), color=GREEN_TEXT, arrow=True)
    add_box_text(slide, "基础题写器\n\n只接收公开事实\n只写清洁题面题意\n记录这道题用了哪些公开事实编号\n只说明题面组织方式\n不解释答案", X(357), Y(153), X(228), Y(126), fill=MINT, line="B4CDBE", size=8.3, color=GREEN_TEXT, margin=2)
    add_box_text(slide, "禁止跨界互动：标准答案/答案推理过程不得影响已明确的目标假设；先生成事实盘点，再让 rubric 评分器重建题面。", X(53), Y(297), X(488), Y(22), fill=PINK, line="E3BDBC", size=6.7, color=RED, margin=1)
    add_page_number(slide, 5)


def build_slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "校验与分流：记录诊断，优先修复，保留探索", "新校验器按 L0-L6 分层输出风险标签，重点用作建议和分流。")
    labels = [
        ("L0", "身份、Schema、旧评分状态"),
        ("L1", "事实来源、规则有效性、场景一致"),
        ("L2", "可答性、事实冲突、题外知识"),
        ("L3", "答案泄漏、单句直达答案"),
        ("L4", "复杂度、竞争存活、信息量平衡"),
        ("L5", "消歧、名称交换、顺序交换"),
        ("L6", "预测结构收益"),
    ]
    for i, (level, desc) in enumerate(labels):
        y = 57 + i * 28
        add_box_text(slide, level, X(53), Y(y), X(34), Y(20), fill=BLUE, line="B6C7CF", size=7.5, color=BLUE_TEXT, margin=0)
        add_text(slide, desc, X(97), Y(y + 1), X(185), Y(20), size=8.2, bold=True, valign=MSO_ANCHOR.MIDDLE)
    ff = slide.shapes.build_freeform(Inches(X(273)), Inches(Y(57)))
    ff.add_line_segments([
        (Inches(X(282)), Inches(Y(57))),
        (Inches(X(282)), Inches(Y(143))),
        (Inches(X(293)), Inches(Y(148))),
        (Inches(X(282)), Inches(Y(153))),
        (Inches(X(282)), Inches(Y(237))),
        (Inches(X(273)), Inches(Y(237))),
    ], close=False)
    bracket = ff.convert_to_shape()
    bracket.fill.background()
    bracket.line.color.rgb = rgb("111111")
    bracket.line.width = Pt(1.5)
    actions = ["只记录风险", "同算子定向重试", "进入预算探索", "人工复核队列", "透传候选"]
    for i, action in enumerate(actions):
        add_box_text(slide, action, X(337), Y(69 + i * 35), X(126), Y(23), fill=GRAY, line="CFD4D4", size=7.6, color=MUTED, margin=1)
    add_box_text(slide, "只有 JSON 无法解析、必填字段缺失、题目为空等技术不可执行问题，才停止候选下游。", X(17), Y(275), X(500), Y(24), fill=YELLOW, line="E1CF8F", size=7.9, color=AMBER, margin=1)
    add_page_number(slide, 6)


def build_slide_7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "评分闭环：新题必须用新答案和新 Rubric", "旧评分材料只可审计，不得作为新评分依据；版本不一致必须停止。")
    labels = ["候选题", "重建参考答案", "生成新 rubric", "重新回答", "Judge打分"]
    xs = [32, 147, 260, 352, 451]
    ws = [95, 92, 71, 78, 74]
    for i, (label, x, w) in enumerate(zip(labels, xs, ws)):
        add_box_text(slide, label, X(x), Y(51), X(w), Y(23), fill=MINT, line="B2D2C0", size=7.7, margin=1)
        if i < len(labels) - 1:
            add_line(slide, X(x + w + 4), Y(63), X(xs[i + 1] - 5), Y(63), color=GREEN_TEXT, arrow=True)
    add_card(slide, "重建原则", ["参考答案只引用公开事实和有效规则", "校验失败则停止在 Rubric 前"], X(32), Y(79), X(151), Y(63), fill="FFFFFF", title_color=GREEN_TEXT, body_size=7.7, title_size=9.0)
    add_card(slide, "回答原则", ["弱模型只看最终题面", "评分标准来自已校验新答案"], X(207), Y(79), X(151), Y(63), fill="FFFFFF", title_color=BLUE_TEXT, body_size=7.7, title_size=9.0)
    add_card(slide, "版本原则", ["题目、答案、Rubric 版本一致", "不一致时拒绝评分"], X(382), Y(79), X(150), Y(63), fill="FFFFFF", title_color=RED, body_size=7.7, title_size=9.0)
    add_section(slide, "效果确认：只把真实能力边界写入经验库", "成功经验必须由稳定降分、目标错误命中和可审计归因共同证明。", title_y=155, sub_y=174, line_y=187)
    rules = [
        "1. 弱模型多次回答后稳定降分",
        "2. 错误命中当前算子目标能力",
        "3. 新参考答案或弱模型答案正确",
        "4. 评分器评分稳定",
        "5. 题面无破坏实验有效性的泄漏或冲突",
        "6. 降分不是格式复杂或评分污染造成",
    ]
    for i, rule in enumerate(rules):
        col = i % 2
        row = i // 2
        x = 37 if col == 0 else 297
        add_box_text(slide, rule, X(x), Y(193 + row * 32), X(234), Y(27), fill=GRAY, line="D1D6D6", size=7.2, align=PP_ALIGN.LEFT, margin=4)
    add_box_text(slide, "当前原则：稳定降分且命中目标错误时保留为边界样本，停止当前分支，不自动纵向叠加同类难度。", X(37), Y(294), X(495), Y(22), fill=MINT, line="B6D1C1", size=7.2, color=GREEN_TEXT, margin=1)
    add_page_number(slide, 7)


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    # Remove the default slide and keep only the blank layout generated pages.
    if prs.slides:
        slide_id = prs.slides._sldIdLst[0]
        prs.part.drop_rel(slide_id.rId)
        prs.slides._sldIdLst.remove(slide_id)
    for builder in (
        build_slide_1,
        build_slide_2,
        build_slide_3,
        build_slide_4,
        build_slide_5,
        build_slide_6,
        build_slide_7,
    ):
        builder(prs)
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    print(build_presentation())
