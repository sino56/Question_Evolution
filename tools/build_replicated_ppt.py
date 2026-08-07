from __future__ import annotations

import argparse
from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUTPUT = ROOT / "Question_Evolution_筛选进化与工程优化架构_无水印.pptx"

REF_W = 592.0
REF_H = 333.0
SLIDE_W = 13.333333
SLIDE_H = 7.5

FONT_CN = "Microsoft YaHei"
FONT_TITLE = "Microsoft YaHei"
FONT_MONO = "Arial"
FONT_SCALE = 1.45


def ix(px: float) -> int:
    return Inches(px / REF_W * SLIDE_W)


def iy(px: float) -> int:
    return Inches(px / REF_H * SLIDE_H)


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def set_cell_border(cell, color: str = "FFFFFF", width: int = 9000) -> None:
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for edge in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        node = tcPr.find(edge, tcPr.nsmap)
        if node is not None:
            tcPr.remove(node)
        tcPr.append(
            parse_xml(
                f'<{edge} w="{width}" {nsdecls("a")}>'
                f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                f'<a:prstDash val="solid"/><a:round/><a:headEnd type="none" w="med" len="med"/>'
                f'<a:tailEnd type="none" w="med" len="med"/></{edge}>'
            )
        )


def set_dash(shape, dash: str = "dash") -> None:
    ln = shape._element.spPr.get_or_add_ln()
    old = ln.find("{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash")
    if old is not None:
        ln.remove(old)
    ln.append(parse_xml(f'<a:prstDash val="{dash}" {nsdecls("a")}/>'))


def no_line(shape) -> None:
    shape.line.fill.background()


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = "FFFFFF",
    line: str | None = None,
    line_width: float = 1.0,
    rounded: bool = False,
    transparency: int = 0,
    dash: bool = False,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, ix(x), iy(y), ix(w), iy(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    if line is None:
        no_line(shape)
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
        if dash:
            set_dash(shape)
    return shape


def set_text_frame(
    frame,
    text: str,
    *,
    size: float = 12,
    color: str = "111111",
    bold: bool = False,
    font: str = FONT_CN,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
    margin: float = 0.8,
    line_spacing: float = 1.0,
) -> None:
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = None
    frame.margin_left = Pt(margin)
    frame.margin_right = Pt(margin)
    frame.margin_top = Pt(margin)
    frame.margin_bottom = Pt(margin)
    frame.vertical_anchor = valign
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = line
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = line_spacing
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size * FONT_SCALE)
            run.font.bold = bold
            run.font.color.rgb = rgb(color)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 12,
    color: str = "111111",
    bold: bool = False,
    font: str = FONT_CN,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
    margin: float = 0.5,
    line_spacing: float = 1.0,
):
    box = slide.shapes.add_textbox(ix(x), iy(y), ix(w), iy(h))
    set_text_frame(
        box.text_frame,
        text,
        size=size,
        color=color,
        bold=bold,
        font=font,
        align=align,
        valign=valign,
        margin=margin,
        line_spacing=line_spacing,
    )
    return box


def add_box_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = "FFFFFF",
    line: str = "222222",
    line_width: float = 1.2,
    rounded: bool = True,
    dash: bool = False,
    size: float = 11,
    color: str = "111111",
    bold: bool = True,
    align: PP_ALIGN = PP_ALIGN.CENTER,
    margin: float = 1.2,
):
    shape = add_rect(
        slide,
        x,
        y,
        w,
        h,
        fill=fill,
        line=line,
        line_width=line_width,
        rounded=rounded,
        dash=dash,
    )
    set_text_frame(
        shape.text_frame,
        text,
        size=size,
        color=color,
        bold=bold,
        align=align,
        margin=margin,
    )
    return shape


def add_line(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: str = "222222",
    width: float = 1.2,
    dash: bool = False,
):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, ix(x1), iy(y1), ix(x2), iy(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if dash:
        set_dash(line)
    return line


def add_arrow(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    color: str = "43BCAF",
    direction: str = "right",
):
    shape_type = {
        "right": MSO_SHAPE.RIGHT_ARROW,
        "left": MSO_SHAPE.LEFT_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
    }[direction]
    shape = slide.shapes.add_shape(shape_type, ix(x), iy(y), ix(w), iy(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    no_line(shape)
    return shape


def add_diamond(slide, text: str, x: float, y: float, w: float, h: float, *, line: str, fill: str = "FFFFFF", size: float = 10):
    shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, ix(x), iy(y), ix(w), iy(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1.2)
    set_text_frame(shape.text_frame, text, size=size, bold=True, align=PP_ALIGN.CENTER, margin=0.3)
    return shape


def add_title(slide, title: str, *, subtitle: str | None = None) -> None:
    add_text(slide, title, 23, 18, 540, 28, size=20, bold=True, valign=MSO_ANCHOR.BOTTOM, margin=0)
    add_line(slide, 23, 49, 568, 49, color="293A3D", width=1.15)
    if subtitle:
        add_text(slide, subtitle, 25, 52, 535, 20, size=8.8, bold=True, valign=MSO_ANCHOR.TOP, margin=0)


def draw_table(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    rows: Sequence[Sequence[str]],
    col_weights: Sequence[float],
    *,
    header_fill: str = "3E86C2",
    row_fills: Sequence[str] = ("C7D8EA", "E1EAF3"),
    header_size: float = 9.5,
    body_size: float = 9.2,
    header_color: str = "FFFFFF",
    alignments: Sequence[PP_ALIGN] | None = None,
    row_heights: Sequence[float] | None = None,
) -> None:
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), ix(x), iy(y), ix(w), iy(h))
    table = table_shape.table
    total = sum(col_weights)
    for idx, weight in enumerate(col_weights):
        table.columns[idx].width = int(ix(w) * weight / total)
    if row_heights:
        total_h = sum(row_heights)
        for idx, weight in enumerate(row_heights):
            table.rows[idx].height = int(iy(h) * weight / total_h)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(header_fill if r == 0 else row_fills[(r - 1) % len(row_fills)])
            align = alignments[c] if alignments else PP_ALIGN.CENTER
            set_text_frame(
                cell.text_frame,
                value,
                size=header_size if r == 0 else body_size,
                color=header_color if r == 0 else "111111",
                bold=True,
                align=align,
                margin=0.5,
                line_spacing=0.9,
            )
            set_cell_border(cell, "F4F7F9", 8500)


def slide1(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = rgb("FFFFFF")
    add_text(s, "样本筛选逻辑", 4, 4, 160, 22, size=15.5, bold=True, margin=0)
    rows = [
        ["维度", "满分", "考察内容"],
        ["主要能力清晰度", "4", "这题有没有一个明确的核心考点，比如“判断前后是否同一人”“判断路线时间是否接得上”"],
        ["进化空间", "4", "后续能不能自然加难，比如加入一个相似人员、监控盲区、时间差、竞争解释。如果只能靠加长题干或罗列更多要点，就不符合。"],
        ["答案/Rubric可重建性", "4", "原答案和评分标准不能支撑后续改写；改写后不能直接复用旧答案，所以原题里必须有清楚的判断规则和边界。"],
        ["题内合成安全性", "3", "加难时是否只需要题内设定，而不是编造真实法条、真实规则、真实案件事实。"],
        ["结果可归因性", "3", "弱模型答错后，能不能知道它主要错在哪里，比如错在“没排除相似目标”，而不是因为题目太散、要求太多。"],
        ["实验信号价值", "2", "运行后是否可能产生有解释价值的升分、降分或不变样本"],
    ]
    draw_table(s, 5, 27, 570, 154, rows, [157, 78, 335], body_size=8.0, row_heights=[18, 38, 43, 42, 27, 31, 25])
    add_text(s, "入选标准", 4, 184, 100, 16, size=14.5, bold=True, margin=0)
    rows2 = [
        ["总分", "结论", ""],
        ["17-20", "优先入选", "判断点清晰，适合进入主实验"],
        ["13-16", "探索入选，根据人工审核是否入选", "有潜力，但风险或信息不足，适合小规模试验"],
        ["9-12", "保留为基线，不入选", "有普通问答价值，但本轮不优先进化"],
        ["0-8", "暂不进入实验", "变难空间弱，或者信息后难以解释原因"],
        ["数据损坏", "需修复", "题目、答案、Rubric 或 JSON 结构有问题"],
    ]
    draw_table(s, 5, 199, 402, 133, rows2, [135, 137, 130], body_size=8.2, row_heights=[18, 21, 29, 25, 29, 24])
    # Screenshot-only ring/page badge are intentionally omitted.


def slide2(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = rgb("FFFFFF")
    add_text(s, "入选样本情况", 4, 4, 185, 23, size=16, bold=True, margin=0)
    add_text(s, "通用场景", 4, 29, 120, 16, size=11.5, bold=True, margin=0)
    add_text(s, "四大特定场景（涉黄、电频车盗窃、拉车门盗窃、夹气运输）", 293, 29, 286, 16, size=10.5, bold=True, margin=0)
    common = [
        ["决策", "分数范围", "样本数", "占比"],
        ["优先入选", "17-20", "218", "72.7%"],
        ["探索入选", "13-16", "48", "16.0%"],
        ["不入选/保留", "9-12", "34", "11.3%"],
    ]
    special = [
        ["决策", "分数范围", "样本数", "占比"],
        ["优先入选", "17-20", "73", "68.9%"],
        ["探索入选", "13-16", "20", "18.9%"],
        ["不入选/保留", "9-12", "13", "12.3%"],
    ]
    draw_table(s, 5, 44, 259, 72, common, [80, 60, 59, 60], body_size=8.5)
    draw_table(s, 294, 44, 258, 72, special, [79, 60, 59, 60], body_size=8.5)
    add_text(s, "注：通用场景考虑数量，仅入选score>=19的，共85条", 5, 118, 280, 14, size=8.5, color="C9272D", bold=True, margin=0)
    add_text(s, "具体类别", 4, 136, 120, 15, size=11.5, bold=True, margin=0)
    categories = [
        ["类别", "说明"],
        ["路线/时间/盲区接续", "判断轨迹是否连续、盲区前后能否接上、是否跟丢"],
        ["多人关系/分工协作", "判断同伙、望风、接应、伴随、分工"],
        ["物品/车辆流转", "判断赃物、气瓶、纸箱、车辆、工具是否一路流转"],
        ["证据能说明到哪一步", "判断视频线索最多能支持什么结论"],
        ["关键动作链条", "判断动作顺序是否完整，比如剪锁、干扰锁车、交付"],
        ["前后是否同一目标", "判断跨衣、遮挡、盲区后是否仍是同一人"],
        ["画面可疑性/行为规律/概念边界", "数据较少，主要覆盖边界、异常规律或操作流程现实"],
    ]
    draw_table(s, 5, 153, 547, 157, categories, [272, 275], body_size=8.0, row_heights=[18, 20, 20, 20, 20, 20, 20, 24])


def slide3(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = rgb("FFFFFF")
    add_text(s, "增量进化", 4, 5, 110, 21, size=17, bold=True, margin=0)
    add_text(s, "流程图：", 4, 25, 85, 15, size=10.5, bold=True, margin=0)

    red = "B43D45"; black = "222222"; light = "F9FAFB"
    add_box_text(s, "输入", 111, 34, 54, 20, line=black, size=9.2)
    add_line(s, 138, 54, 138, 72, width=1.2)
    add_box_text(s, "初始问题样本", 99, 72, 77, 26, line=black, size=8.5)
    add_arrow(s, 178, 79, 15, 10, color="C9C9C9")
    add_box_text(s, "评分与样本分类\n策略", 194, 71, 79, 29, line=black, size=8.0)
    add_line(s, 138, 98, 138, 111, width=1.2)
    add_box_text(s, "题目 / 答案 / rubric\n写入 JSON", 91, 111, 94, 37, line=black, size=7.7)

    large = add_rect(s, 7, 157, 257, 82, fill="FFFFFF", line=red, line_width=1.1, rounded=True, dash=True)
    add_text(s, "样本合法性检查", 16, 190, 69, 22, size=8.8, bold=True, align=PP_ALIGN.CENTER)
    add_line(s, 82, 166, 82, 225, width=1.1)
    add_text(
        s,
        "1. 校验字段完整、题面和答案可解析；\n2. 根据规则检查重复、冲突或异常内容；\n3. 结合评分与边界信号判断：\n   · 满足条件 → 进入后续进化；\n   · 存在硬风险 → 跳过；\n   · 边界样本 → 记录失败原因。",
        87, 162, 167, 66, size=7.4, bold=True, margin=0, line_spacing=0.9,
    )
    add_line(s, 264, 198, 302, 198, width=1.2)

    top = add_rect(s, 263, 14, 143, 42, fill="FFFFFF", line=red, line_width=1.0, rounded=True, dash=True)
    add_text(s, "1. 读取进化策略与历史失败信号；\n2. 选择候选 operator 并生成改写；\n3. 保存进化问题与变更原因。", 270, 18, 128, 34, size=7.1, bold=True, margin=0)
    add_line(s, 334, 56, 334, 88, width=1.3)
    add_box_text(s, "gpm=object-memory\n选择更合适的改写方向\n避免重复失败模式", 281, 89, 105, 60, line=black, size=8.0)
    add_line(s, 334, 149, 334, 157, width=1.2)
    add_diamond(s, "是否进化?", 302, 157, 67, 61, line=black, size=9.5)
    add_text(s, "否", 372, 188, 18, 14, size=8.5, bold=True, margin=0)
    add_line(s, 369, 188, 433, 188, color="999999", dash=True)
    add_text(s, "是", 337, 143, 18, 14, size=8.5, bold=True, margin=0)

    add_text(s, "Round_N+1", 268, 224, 70, 15, size=8.3, font=FONT_MONO, bold=True, margin=0)
    add_line(s, 334, 218, 334, 233, width=1.2)
    add_box_text(s, "更新样本状态\n保存失败原因", 342, 234, 91, 34, line=black, size=7.4)
    add_text(s, "否", 324, 230, 16, 14, size=8.3, bold=True, margin=0)
    add_line(s, 334, 268, 334, 286, width=1.2)
    add_text(s, "停止", 308, 282, 30, 15, size=8.0, bold=True, margin=0)
    add_text(s, "是", 354, 282, 22, 15, size=8.0, bold=True, margin=0)

    add_box_text(s, "生成候选问题", 462, 44, 70, 27, line=black, size=8.2)
    add_line(s, 497, 71, 497, 82, width=1.2)
    add_box_text(s, "gpt生成候选答案、rubric、\nanswer / score result", 433, 81, 124, 46, line=black, size=7.3)
    right_group = add_rect(s, 454, 130, 102, 73, fill="FFFFFF", line=red, line_width=1.0, rounded=True, dash=True)
    add_box_text(s, "Qwen候选回答\n（多次采样）", 463, 135, 83, 25, line=black, size=7.4)
    add_box_text(s, "聚合评分 / 稳定性\n均值与方差", 463, 165, 83, 25, line=black, size=7.2)
    add_line(s, 504, 127, 504, 135, width=1.2)
    add_line(s, 504, 160, 504, 165, width=1.2)
    add_line(s, 504, 190, 504, 205, width=1.2)
    add_box_text(s, "选择最优候选\n进入主链评分", 462, 205, 87, 32, line=black, size=7.3)
    add_line(s, 505, 237, 505, 251, width=1.2)
    note = add_rect(s, 435, 251, 124, 38, fill="FFFFFF", line=red, line_width=1.0, rounded=True, dash=True)
    add_text(s, "1. 对比进化前后得分变化；\n2. 判断是否形成有效边界；\n3. 记录成功与失败模式。", 441, 255, 112, 30, size=7.0, bold=True, margin=0)
    add_text(s, "校验失败时不产出，导致样本无法进入进化流程", 52, 258, 215, 18, size=8.5, color=red, bold=True, margin=0)
    add_box_text(s, "成功：记录边界与 operator 收益\n失败：写入 failure memory\n无增益：切换策略继续探索\n不可用：停止或进入人工复核", 346, 290, 244, 42, fill="FFF8F8", line=red, line_width=1.0, dash=True, size=7.0, align=PP_ALIGN.LEFT)


def slide4(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = rgb("FFFFFF")
    add_title(s, "后续树搜索+回溯+算子叠加结构")
    add_rect(s, 4, 50, 584, 76, fill="D7E7F4", rounded=True)
    add_rect(s, 4, 130, 584, 102, fill="FFF0B9", rounded=True)
    add_rect(s, 4, 236, 584, 94, fill="DDEFD5", rounded=True)
    add_text(s, "第0层", 19, 69, 60, 32, size=16, bold=True, margin=0)
    add_text(s, "第1层", 19, 170, 60, 32, size=16, bold=True, margin=0)
    add_text(s, "第2层", 19, 273, 60, 32, size=16, bold=True, margin=0)
    add_box_text(s, "ROOT", 255, 53, 82, 23, fill="EEF7FE", line="111111", size=10.5)
    add_line(s, 296, 76, 296, 86, width=1.4)
    add_box_text(s, "候选算子列表", 232, 86, 129, 29, fill="FBE0CF", line="111111", size=10.5)
    add_box_text(s, "扩展最大数量限制", 232, 116, 129, 17, fill="F7E7E7", line="C62E32", line_width=1.0, dash=True, size=9, color="C62E32")
    add_box_text(s, "候选1", 126, 139, 57, 25, fill="FFF8D6", line="111111", size=9.5)
    add_box_text(s, "候选2", 224, 139, 57, 25, fill="FFF8D6", line="111111", size=9.5)
    add_text(s, "…", 299, 142, 25, 20, size=12, bold=True, align=PP_ALIGN.CENTER)
    add_box_text(s, "候选N", 329, 139, 61, 25, fill="FFF8D6", line="111111", size=9.5, color="D82024")
    add_line(s, 296, 133, 154, 139, width=1.1)
    add_line(s, 296, 133, 252, 139, width=1.1)
    add_line(s, 296, 133, 360, 139, color="C62E32", width=1.1)
    add_box_text(s, "叠加算子候选列表", 90, 197, 128, 28, fill="FBE8DA", line="111111", size=9.5)
    add_line(s, 154, 164, 154, 197, width=1.2)
    add_text(s, "有效\n进化", 157, 163, 32, 32, size=9, color="D42025", bold=True, align=PP_ALIGN.CENTER, margin=0)
    add_line(s, 252, 164, 252, 196, width=1.2)
    add_text(s, "无效\n进化", 247, 166, 37, 31, size=9, color="D42025", bold=True, align=PP_ALIGN.CENTER, margin=0)
    add_text(s, "回溯", 207, 169, 32, 19, size=9, color="D42025", bold=True, margin=0)
    back = add_line(s, 222, 216, 222, 121, color="C62E32", width=1.1, dash=True)
    add_box_text(s, "候选1+A1", 71, 255, 74, 27, fill="EAF7E5", line="111111", size=9.2)
    add_box_text(s, "候选1+An", 176, 255, 77, 27, fill="EAF7E5", line="111111", size=9.2)
    add_text(s, "…", 151, 260, 23, 18, size=11, bold=True, align=PP_ALIGN.CENTER)
    add_line(s, 154, 225, 108, 255, width=1.1)
    add_line(s, 154, 225, 214, 255, width=1.1)
    add_line(s, 108, 282, 108, 309, width=1.1)
    add_line(s, 214, 282, 214, 309, width=1.1)
    add_text(s, "有效\n进化", 111, 285, 34, 29, size=8.8, color="D42025", bold=True, align=PP_ALIGN.CENTER, margin=0)
    add_text(s, "无效\n进化", 219, 285, 34, 29, size=8.8, color="D42025", bold=True, align=PP_ALIGN.CENTER, margin=0)
    add_text(s, "记录", 93, 311, 40, 17, size=9, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "停止", 200, 311, 40, 17, size=9, bold=True, align=PP_ALIGN.CENTER)
    callout = add_box_text(s, "全局停止条件：所有候选算子都已遍历或者每个样本边界数达到最大数量限制", 407, 208, 147, 55, fill="FFFFFF", line="111111", size=8.8, align=PP_ALIGN.LEFT)
    # Emphasize the prefix in red with a short overlay.
    add_text(s, "全局停止条件：", 417, 216, 84, 16, size=9, color="D42025", bold=True, margin=0)


def slide5(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = rgb("FFFFFF")
    add_title(s, "评分阶段：受控并行架构", subtitle="样本 worker、Qwen 回答池和 GPT Judge 调度相互隔离；并发上限约为 20，Qwen 回答和 Judge 共享同一套限流。")
    add_box_text(s, "评分样本 Worker × N（有界）", 218, 76, 178, 25, fill="FFFFFF", line="8E77CE", size=9.7)
    add_arrow(s, 297, 102, 20, 12, color="9678DF", direction="down")
    add_box_text(s, "Qwen 回答 Trial × n", 112, 120, 169, 23, fill="FFFFFF", line="A4DAD3", size=8.7)
    add_box_text(s, "GPT Judge Repeat × m", 335, 120, 168, 23, fill="FFFFFF", line="E7D59C", size=8.7)
    add_arrow(s, 184, 144, 17, 11, color="4BC5C6", direction="down")
    add_arrow(s, 409, 144, 17, 11, color="F4C34B", direction="down")
    add_box_text(s, "Qwen 公平调度池\n独立限流 · 上限 20", 88, 157, 208, 34, fill="FFFFFF", line="2CB5C3", line_width=1.4, size=8.5)
    add_box_text(s, "GPT 公平调度池\n独立限流 · 上限 20", 325, 157, 208, 34, fill="FFFFFF", line="D9A700", line_width=1.4, size=8.5)
    add_box_text(s, "优化前\n串行：回答一次，测评\n样本并发 = 服务并发", 24, 203, 108, 49, fill="FFFFFF", line="E3E3E3", size=8.1, align=PP_ALIGN.CENTER)
    add_box_text(s, "公平调度策略\n多活跃样本 → 优先让不同样本各获一个名额\n有富余 → 再处理同一样本的剩余 trial / repeat", 141, 201, 296, 52, fill="FFFFFF", line="B8DED7", line_width=0.9, dash=True, size=8.0)
    add_box_text(s, "优化后\n样本内 trial 并行\n服务并发受总上限约束", 450, 203, 117, 49, fill="FFFFFF", line="C9E2DF", size=8.1)
    add_line(s, 20, 260, 568, 260, color="E4E4E4", width=0.8)
    add_text(s, "样本消费：每个样本参与固定在公平调度队列；", 24, 269, 320, 15, size=8.0, bold=True, margin=0)
    add_text(s, "回答消费：每一轮都只在独立限流器内消费，避免挤占服务资源；", 24, 286, 440, 15, size=8.0, bold=True, margin=0)
    add_text(s, "最终收益：样本内 trial 并行 + 跨样本公平分配，将『样本并发 = 服务并发』解耦为『样本并发 × 服务并发』，服务并发受独立总上限约束。", 24, 306, 548, 21, size=7.7, bold=True, margin=0)


def slide6(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = rgb("FFFFFF")
    add_title(s, "有界任务调度与流式 I/O", subtitle="Reader 仅在队列有容量时提交记录，Writer 批量聚合、内存、协程和落盘速率被限制在并发配置附近，不随输入线性增长。")
    add_box_text(s, "流式 Reader\n逐行读取 JSONL", 24, 115, 83, 43, fill="FFFFFF", line="3FC1C5", size=8.2)
    add_arrow(s, 109, 126, 22, 17, color="40C4C5")
    add_box_text(s, "输入队列\n有界容量\n背压控制", 130, 112, 77, 49, fill="FFFFFF", line="E4BC40", size=8.0)
    add_arrow(s, 210, 126, 20, 17, color="F0BE43")
    pool = add_rect(s, 234, 95, 106, 90, fill="FFFFFF", line="8176BE", line_width=1.0, rounded=True, dash=True)
    add_text(s, "Worker 池（固定）", 247, 98, 82, 16, size=8.4, bold=True, align=PP_ALIGN.CENTER, margin=0)
    add_box_text(s, "Worker 1", 243, 116, 88, 18, fill="FFFFFF", line="B0A7D2", line_width=0.8, size=7.6)
    add_box_text(s, "Worker 2", 243, 137, 88, 18, fill="FFFFFF", line="B0A7D2", line_width=0.8, size=7.6)
    add_box_text(s, "Worker ... N", 243, 158, 88, 18, fill="FFFFFF", line="B0A7D2", line_width=0.8, size=7.6)
    add_arrow(s, 343, 126, 20, 17, color="44BDA5")
    add_box_text(s, "输出队列\n有界容量", 364, 112, 71, 49, fill="FFFFFF", line="48BBA5", size=8.0)
    add_arrow(s, 437, 126, 20, 17, color="44BDA5")
    add_box_text(s, "Batch Writer\n• 按数量、字节或时间阈值批量 flush，减少频繁磁盘同步；\n• 避免多 worker 并发写文件造成乱序或竞争", 458, 94, 105, 86, fill="FFFFFF", line="3C8BC0", size=7.8, align=PP_ALIGN.LEFT)
    add_line(s, 168, 162, 168, 193, color="E8C65C", width=0.9)
    add_line(s, 168, 193, 399, 193, color="E8C65C", width=0.9)
    add_line(s, 399, 193, 399, 162, color="E8C65C", width=0.9)
    add_text(s, "背压反馈：队列满时 Reader 暂停提交", 206, 196, 193, 16, size=7.8, bold=True, align=PP_ALIGN.CENTER, margin=0)
    add_box_text(s, "流式读取\n逐行读取 JSONL，跳过空行，即时投递，避免全量加载导致内存峰值。", 24, 244, 254, 51, fill="FFFFFF", line="BDE5E1", line_width=0.8, size=8.0, align=PP_ALIGN.LEFT)
    add_box_text(s, "批量写入 + partial文件\nWriter 按记录数/字节/时间阈值批量 flush，阶段完成结果 → 同步写入checkpoint便于后续中断恢复 → 降低内存峰值。", 291, 244, 276, 51, fill="FFFFFF", line="D0E3E1", line_width=0.8, size=7.8, align=PP_ALIGN.LEFT)


def slide7(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = rgb("FFFFFF")
    add_title(s, "原子发布与中断恢复", subtitle="正式产物只在阶段完整成功后才可见，中断后可从 checkpoint 恢复，避免重复计算。")
    add_box_text(s, "业务 Worker\n产生结果", 23, 110, 91, 35, fill="FFFFFF", line="8179C4", size=8.4)
    add_arrow(s, 116, 119, 21, 16, color="8972D5")
    add_box_text(s, "Partial 文件\n尚未发布的记录", 135, 108, 91, 39, fill="FFFFFF", line="DA7496", size=8.2)
    add_arrow(s, 228, 119, 21, 16, color="E36F9A")
    add_box_text(s, "Checkpoint\n已持久化位置", 248, 108, 89, 39, fill="FFFFFF", line="DA7496", size=8.2)
    add_arrow(s, 339, 119, 21, 16, color="E36F9A")
    add_diamond(s, "阶段\n完成?", 359, 107, 50, 41, line="D4AF34", fill="FFFDF5", size=8.0)
    add_text(s, "中断", 410, 102, 34, 16, size=8.0, bold=True, margin=0)
    add_arrow(s, 412, 119, 20, 16, color="E36F78")
    add_box_text(s, "重启恢复 -\n读取 checkpoint", 430, 106, 87, 42, fill="FFFFFF", line="D27780", dash=True, size=8.0)
    add_text(s, "完成", 393, 150, 31, 15, size=8.0, bold=True, margin=0)
    add_arrow(s, 393, 149, 13, 18, color="46BEA2", direction="down")
    add_box_text(s, "原子链接发布\n正式产物（重命名临时文件）\nManifest（校验信息文件）", 296, 166, 120, 47, fill="FFFFFF", line="3BB89D", size=7.1)
    add_arrow(s, 476, 149, 13, 18, color="E67790", direction="down")
    add_box_text(s, "跳过已完成记录\n继续未完成部分", 425, 166, 95, 47, fill="FFFFFF", line="E4C0C7", size=7.6)
    add_box_text(s, "Manifest 校验（最终值完整）：记录数 + 内容摘要 + 代码版本 + 校验结果 → 下游只看到完整结果", 23, 220, 544, 30, fill="FFFFFF", line="BDE2DE", line_width=0.8, size=7.7)
    add_box_text(s, "✓ 原子发布\n\n正式产物只在阶段完整成功后才可见，避免下游读到一半写入结果。Manifest 已记录条数、内容摘要和代码版本。", 23, 263, 269, 63, fill="FFFFFF", line="BDE2DE", line_width=0.8, size=8.0, align=PP_ALIGN.LEFT)
    add_text(s, "✓", 35, 267, 16, 16, size=10, color="24AF91", bold=True, margin=0)
    add_box_text(s, "↻ 中断恢复\n\npartial + checkpoint 记录已确认进度。中断重启时跳过已完成记录/批次，从断点继续，避免整轮重跑。", 302, 263, 265, 63, fill="FFFFFF", line="E1D7DA", line_width=0.8, size=8.0, align=PP_ALIGN.LEFT)
    add_text(s, "↻", 315, 267, 16, 16, size=10, color="D94973", bold=True, margin=0)


def slide8(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = rgb("FFFFFF")
    add_title(s, "Memory 索引与本地校验复用", subtitle="倒排索引 + 匹配缓存避免全量扫描 memory bank；观察版本一致时复用进化阶段校验结果。")
    add_box_text(s, "Memory Bank\n正 / 负 memory\n按校验增长", 24, 82, 107, 51, fill="FFFFFF", line="42C6AD", size=8.3)
    add_arrow(s, 134, 98, 20, 16, color="43C0A5")
    add_box_text(s, "倒排索引\n基于 signature 构建", 155, 82, 112, 51, fill="FFFFFF", line="3BAFC0", size=8.3)
    add_arrow(s, 270, 98, 20, 16, color="3FBAD1")
    add_box_text(s, "匹配缓存\n相同 signature\n直接命中缓存", 289, 82, 110, 51, fill="FFFFFF", line="7F6BC4", size=8.2)
    add_arrow(s, 402, 98, 20, 16, color="8B66D2")
    add_box_text(s, "输出结果\nprimary / backup\navoid + 理由", 423, 82, 109, 51, fill="FFFFFF", line="D6A628", size=8.2)
    b = add_rect(s, 24, 141, 231, 83, fill="FFFFFF", line="A8D9D2", line_width=0.8, rounded=True, dash=True)
    add_text(s, "每个样本先提取一个 signature，目前主要包含：\n\n• 核心能力\n• 主题场景\n• 问题形态\n• 高分原因/错误原因", 30, 146, 218, 73, size=8.1, bold=True, margin=0, line_spacing=0.95)
    add_arrow(s, 272, 176, 25, 16, color="3CC1C5")
    add_text(s, "建立“字段值 → 包含该值的 memory 记录”的索引，提高检索效率", 307, 169, 247, 28, size=8.2, bold=True, margin=0)
    add_box_text(s, "◆ 索引复用\n\n基于 signature 构建只读倒排索引，加载一次。匹配缓存避免重复计算相同 signature 的路由结果。", 24, 239, 267, 57, fill="FFFFFF", line="BDE2DE", line_width=0.8, size=8.0, align=PP_ALIGN.LEFT)
    add_text(s, "◆", 35, 244, 16, 14, size=10, color="19AAC8", bold=True, margin=0)
    add_box_text(s, "◆ 校验复用\n\n同源同版本直接复用进化阶段校验结果，跳过重复评估；LLM 校验不受影响，保证语义安全。", 301, 239, 266, 57, fill="FFFFFF", line="D1E1DF", line_width=0.8, size=8.0, align=PP_ALIGN.LEFT)
    add_text(s, "◆", 312, 244, 16, 14, size=10, color="1EB884", bold=True, margin=0)
    add_box_text(s, "优化后完整运行时间约为原来的40%-60%", 199, 306, 195, 22, fill="FFFFFF", line="222222", line_width=1.1, size=8.6, color="D2262D")


def build(output: Path = DEFAULT_OUTPUT) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    # Remove the starter slide if the template contains one.
    while prs.slides:
        slide_id = prs.slides._sldIdLst[0]
        prs.part.drop_rel(slide_id.rId)
        prs.slides._sldIdLst.remove(slide_id)
    for maker in (slide1, slide2, slide3, slide4, slide5, slide6, slide7, slide8):
        maker(prs)
    output = output.resolve()
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(f"saved: {output}")
    print(f"slides: {len(prs.slides)}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Build the watermark-free editable PPT recreation.")
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    args = parser.parse_args()
    build(args.output)
